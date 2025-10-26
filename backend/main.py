import os
import re
import time
import zipfile
import tarfile
import asyncio
import mimetypes
import csv
import io
from datetime import datetime
from pathlib import Path
from typing import List, Optional, Dict, AsyncGenerator
from contextlib import asynccontextmanager

from fastapi import FastAPI, HTTPException, WebSocket, WebSocketDisconnect
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import StreamingResponse
from sqlalchemy import create_engine, Column, Integer, String, Text, Boolean, BigInteger, DateTime, ForeignKey
from sqlalchemy.ext.declarative import declarative_base
from sqlalchemy.orm import sessionmaker, Session, relationship
from pydantic import BaseModel
import uvicorn

# Document processing libraries
try:
    from docx import Document as DocxDocument
except ImportError:
    DocxDocument = None

try:
    from openpyxl import load_workbook
except ImportError:
    load_workbook = None

try:
    from pptx import Presentation
except ImportError:
    Presentation = None

try:
    from odf.opendocument import load as odf_load
    from odf.text import P as OdfP
except ImportError:
    odf_load = None
    OdfP = None

try:
    from PyPDF2 import PdfReader
except ImportError:
    PdfReader = None


# Security: Path validation function
def validate_and_resolve_path(user_path: str, base_path: Path) -> str:
    """
    Securely validate and resolve a user-provided path.
    Follows OWASP guidelines and CodeQL recommendations for path traversal prevention.
    
    Args:
        user_path: User-provided path string
        base_path: Base directory that all paths must be within
        
    Returns:
        Validated normalized path string that is guaranteed to be within base_path
        
    Raises:
        HTTPException: If path is invalid or outside base_path
    """
    # Convert base_path to string for os.path operations
    base_path_str = str(base_path.resolve())
    
    # Sanitize user input
    # Remove any null bytes
    if '\0' in user_path:
        raise HTTPException(
            status_code=400,
            detail="Invalid path: null bytes not allowed"
        )
    
    # Handle paths that already include the base
    if user_path.startswith("/app/host_root"):
        # Use the path as-is but still validate
        candidate_path_str = user_path
    else:
        # Remove leading slash for joining
        clean_path = user_path.lstrip("/")
        
        # Explicit check for path traversal patterns
        # Check for .. in any part of the path
        path_parts = clean_path.split(os.sep)
        if ".." in path_parts or any(".." in part for part in path_parts):
            raise HTTPException(
                status_code=400,
                detail="Invalid path: path traversal not allowed"
            )
        
        # Check for tilde expansion
        if "~" in clean_path:
            raise HTTPException(
                status_code=400,
                detail="Invalid path: tilde expansion not allowed"
            )
        
        # Check for absolute paths trying to escape
        if os.path.isabs(clean_path):
            raise HTTPException(
                status_code=400,
                detail="Invalid path: absolute paths not allowed"
            )
        
        # Join with base path using os.path.join
        candidate_path_str = os.path.join(base_path_str, clean_path)
    
    # Normalize the path (resolves .., ., // etc)
    # This is the key CodeQL-recommended pattern
    normalized_path = os.path.normpath(candidate_path_str)
    
    # Verify normalized path starts with base_path (prefix check)
    # This is the critical security check
    if not normalized_path.startswith(base_path_str + os.sep) and normalized_path != base_path_str:
        raise HTTPException(
            status_code=403,
            detail="Access denied: Path outside allowed directory"
        )
    
    # Additional symlink check using os.path.realpath (string-based)
    try:
        real_path = os.path.realpath(normalized_path)
        real_base = os.path.realpath(base_path_str)
        
        # Verify the real path is still within base
        if not real_path.startswith(real_base + os.sep) and real_path != real_base:
            raise HTTPException(
                status_code=403,
                detail="Access denied: Resolved path outside allowed directory"
            )
    except (OSError, RuntimeError) as e:
        # Path might not exist yet, which is okay for some operations
        # The prefix check above is the critical security boundary
        pass
    
    # Return the validated string path
    return normalized_path


# Database setup
DATABASE_URL = os.getenv("DATABASE_URL", "postgresql://postgres:password@localhost:5432/filesearch")
engine = create_engine(DATABASE_URL)
SessionLocal = sessionmaker(autocommit=False, autoflush=False, bind=engine)
Base = declarative_base()

# Global cancellation tracking for directory analysis
analysis_cancellations = {}
analysis_cancellations_lock = asyncio.Lock()

# Models
class SearchSession(Base):
    __tablename__ = "search_sessions"
    
    id = Column(Integer, primary_key=True, index=True)
    search_term = Column(String(500), nullable=False)
    search_path = Column(String(1000), nullable=False)
    start_time = Column(DateTime, default=datetime.utcnow)
    end_time = Column(DateTime)
    total_files_searched = Column(Integer, default=0)
    total_matches = Column(Integer, default=0)
    status = Column(String(50), default="running")
    
    results = relationship("SearchResult", back_populates="session", cascade="all, delete-orphan")

class SearchResult(Base):
    __tablename__ = "search_results"
    
    id = Column(Integer, primary_key=True, index=True)
    session_id = Column(Integer, ForeignKey("search_sessions.id"))
    file_path = Column(String(2000), nullable=False)
    file_name = Column(String(500), nullable=False)
    file_size = Column(BigInteger)
    file_type = Column(String(50))
    match_count = Column(Integer, default=0)
    is_zip_file = Column(Boolean, default=False)
    zip_parent_path = Column(String(2000))
    found_at = Column(DateTime, default=datetime.utcnow)
    preview_text = Column(Text)
    
    session = relationship("SearchSession", back_populates="results")
    matches = relationship("MatchDetail", back_populates="result", cascade="all, delete-orphan")

class MatchDetail(Base):
    __tablename__ = "match_details"
    
    id = Column(Integer, primary_key=True, index=True)
    result_id = Column(Integer, ForeignKey("search_results.id"))
    line_number = Column(Integer)
    line_content = Column(Text)
    match_position = Column(Integer)
    context_before = Column(Text)
    context_after = Column(Text)
    
    result = relationship("SearchResult", back_populates="matches")

class SavedSearch(Base):
    __tablename__ = "saved_searches"
    
    id = Column(Integer, primary_key=True, index=True)
    name = Column(String(200), nullable=False, unique=True)
    search_term = Column(String(500), nullable=False)
    search_path = Column(String(1000), nullable=False)
    case_sensitive = Column(Boolean, default=False)
    include_zip_files = Column(Boolean, default=True)
    search_filenames = Column(Boolean, default=False)
    file_type = Column(String(100))  # Optional file type filter
    min_size = Column(BigInteger)  # Optional minimum file size
    max_size = Column(BigInteger)  # Optional maximum file size
    created_at = Column(DateTime, default=datetime.utcnow)

# Pydantic models
class SearchRequest(BaseModel):
    search_term: str
    search_path: str = "/Users/vndr"
    case_sensitive: bool = False
    include_zip_files: bool = True
    search_filenames: bool = False  # New: search by filename instead of content

class SearchSessionResponse(BaseModel):
    id: int
    search_term: str
    search_path: str
    start_time: datetime
    end_time: Optional[datetime]
    total_files_searched: int
    total_matches: int
    status: str

class SearchResultResponse(BaseModel):
    id: int
    file_path: str
    file_name: str
    file_size: Optional[int]
    file_type: Optional[str]
    match_count: int
    is_zip_file: bool
    zip_parent_path: Optional[str]
    preview_text: Optional[str]

class MatchDetailResponse(BaseModel):
    line_number: Optional[int]
    line_content: str
    match_position: int
    context_before: Optional[str]
    context_after: Optional[str]

class SavedSearchCreate(BaseModel):
    name: str
    search_term: str
    search_path: str = "/Users/vndr"
    case_sensitive: bool = False
    include_zip_files: bool = True
    search_filenames: bool = False
    file_type: Optional[str] = None
    min_size: Optional[int] = None
    max_size: Optional[int] = None

class SavedSearchResponse(BaseModel):
    id: int
    name: str
    search_term: str
    search_path: str
    case_sensitive: bool
    include_zip_files: bool
    search_filenames: bool
    file_type: Optional[str]
    min_size: Optional[int]
    max_size: Optional[int]
    created_at: datetime

# File search engine
class FileSearchEngine:
    def __init__(self, db_session: Session):
        self.db_session = db_session
        self.supported_text_extensions = {
            '.txt', '.py', '.js', '.ts', '.html', '.css', '.json', '.xml', '.yml', '.yaml',
            '.md', '.rst', '.log', '.cfg', '.conf', '.ini', '.properties', '.sql',
            '.sh', '.bash', '.zsh', '.ps1', '.bat', '.cmd', '.dockerfile', '.makefile',
            '.cpp', '.c', '.h', '.hpp', '.java', '.cs', '.php', '.rb', '.go', '.rs',
            '.swift', '.kt', '.scala', '.clj', '.pl', '.r', '.m', '.mm', '.lua', '.vim'
        }
    
    async def search_files(self, request: SearchRequest, websocket: Optional[WebSocket] = None) -> int:
        """Main search function that returns session_id"""
        # Create search session
        session = SearchSession(
            search_term=request.search_term,
            search_path=request.search_path,
            status="running"
        )
        self.db_session.add(session)
        self.db_session.commit()
        self.db_session.refresh(session)
        
        session_id = session.id
        start_time = time.time()
        
        try:
            # Create search pattern based on search mode
            if request.search_filenames:
                search_pattern = self._create_filename_pattern(request.search_term, request.case_sensitive)
            else:
                search_pattern = self._create_search_pattern(request.search_term, request.case_sensitive)
            
            # Secure path handling - prevent directory traversal attacks
            base_path = "/app/host_root"
            local_path = request.search_path
            
            # Normalize and join paths
            full_path = os.path.normpath(os.path.join(base_path, local_path.lstrip('/')))
            
            # Resolve to absolute path and verify it's within base_path
            real_base_path = os.path.realpath(base_path)
            real_full_path = os.path.realpath(full_path)
            
            # Security check: ensure the resolved path is within the allowed base directory
            if not real_full_path.startswith(real_base_path + os.sep) and real_full_path != real_base_path:
                raise HTTPException(status_code=400, detail="Access to the requested path is not allowed.")
            
            search_path = Path(real_full_path)
            
            if not search_path.exists():
                raise HTTPException(status_code=400, detail=f"Search path does not exist: {local_path}")
            
            total_files = 0
            total_matches = 0
            
            async for file_path in self._walk_directory(search_path):
                if websocket:
                    try:
                        await websocket.send_json({
                            "type": "progress",
                            "current_file": str(file_path).replace('/app/host_root', ''),  # Show original local path
                            "files_searched": total_files
                        })
                    except Exception:  # nosec B110 - Catch all websocket errors
                        break  # WebSocket disconnected
                
                total_files += 1
                
                try:
                    # If searching by filename, check filename match
                    if request.search_filenames:
                        if search_pattern.search(file_path.name):
                            # Create a result for filename match
                            result = SearchResult(
                                session_id=session_id,
                                file_path=str(file_path).replace('/app/host_root', ''),
                                file_name=file_path.name,
                                file_size=file_path.stat().st_size if file_path.is_file() else 0,
                                file_type=file_path.suffix.lower(),
                                match_count=1,
                                is_zip_file=False,
                                preview_text=f"Filename match: {file_path.name}"
                            )
                            self.db_session.add(result)
                            self.db_session.flush()
                            
                            # Add a single match detail for the filename
                            match_detail = MatchDetail(
                                result_id=result.id,
                                line_number=0,
                                line_content=f"Filename: {file_path.name}",
                                match_position=0
                            )
                            self.db_session.add(match_detail)
                            total_matches += 1
                        continue  # Skip content search when searching filenames
                    
                    # Content search (original logic)
                    file_ext = file_path.suffix.lower()
                    
                    # Handle archive files
                    if file_ext == '.zip' and request.include_zip_files:
                        matches = await self._search_zip_file(file_path, search_pattern, session_id)
                        total_matches += matches
                    elif file_ext in {'.tar', '.gz', '.bz2', '.tgz', '.tar.gz', '.tar.bz2'} and request.include_zip_files:
                        matches = await self._search_tar_file(file_path, search_pattern, session_id)
                        total_matches += matches
                    # Handle Microsoft Office documents
                    elif file_ext in {'.docx', '.doc'}:
                        matches = await self._search_docx_file(file_path, search_pattern, session_id)
                        total_matches += matches
                    elif file_ext in {'.xlsx', '.xls'}:
                        matches = await self._search_xlsx_file(file_path, search_pattern, session_id)
                        total_matches += matches
                    elif file_ext in {'.pptx', '.ppt'}:
                        matches = await self._search_pptx_file(file_path, search_pattern, session_id)
                        total_matches += matches
                    # Handle LibreOffice/OpenDocument files
                    elif file_ext in {'.odt', '.ods', '.odp'}:
                        matches = await self._search_odf_file(file_path, search_pattern, session_id)
                        total_matches += matches
                    # Handle regular text files
                    elif self._is_text_file(file_path):
                        matches = await self._search_text_file(file_path, search_pattern, session_id)
                        total_matches += matches
                        
                except Exception as e:
                    print(f"Error searching file {file_path}: {e}")
                    continue
                
                # Update progress every 100 files
                if total_files % 100 == 0:
                    session.total_files_searched = total_files
                    session.total_matches = total_matches
                    self.db_session.commit()
            
            # Finalize session
            end_time = time.time()
            session.total_files_searched = total_files
            session.total_matches = total_matches
            session.end_time = datetime.utcnow()
            session.status = "completed"
            self.db_session.commit()
            
            if websocket:
                try:
                    await websocket.send_json({
                        "type": "completed",
                        "session_id": session_id,
                        "total_files": total_files,
                        "total_matches": total_matches,
                        "duration": end_time - start_time
                    })
                except Exception:  # nosec B110 - Catch all websocket errors
                    pass
            
            return session_id
            
        except Exception as e:
            session.status = "error"
            self.db_session.commit()
            raise e
    
    def _create_search_pattern(self, search_term: str, case_sensitive: bool):
        """Create pattern for content search (exact match)"""
        flags = 0 if case_sensitive else re.IGNORECASE
        return re.compile(re.escape(search_term), flags)
    
    def _create_filename_pattern(self, search_term: str, case_sensitive: bool):
        """Create pattern for filename search (supports wildcards and partial match)"""
        flags = 0 if case_sensitive else re.IGNORECASE
        
        # Convert shell-style wildcards to regex
        # * becomes .* (match any characters)
        # ? becomes . (match single character)
        pattern = search_term
        
        # Check if user is using wildcards
        if '*' in pattern or '?' in pattern:
            # Escape special regex characters except * and ?
            pattern = re.escape(pattern)
            # Convert shell wildcards to regex
            pattern = pattern.replace(r'\*', '.*').replace(r'\?', '.')
        else:
            # No wildcards - do partial match (search anywhere in filename)
            pattern = re.escape(pattern)
        
        # Compile with appropriate flags
        return re.compile(pattern, flags)
    
    async def _walk_directory(self, path: Path) -> AsyncGenerator[Path, None]:
        """Async generator to walk through directory"""
        for root, dirs, files in os.walk(path):
            # Skip hidden directories and common non-searchable directories
            dirs[:] = [d for d in dirs if not d.startswith('.') and d not in {'node_modules', '__pycache__', '.git'}]
            
            for file in files:
                if not file.startswith('.'):
                    yield Path(root) / file
    
    def _is_text_file(self, file_path: Path) -> bool:
        """Check if file is likely to be a text file"""
        if file_path.suffix.lower() in self.supported_text_extensions:
            return True
        
        # Check MIME type
        mime_type, _ = mimetypes.guess_type(str(file_path))
        if mime_type and mime_type.startswith('text'):
            return True
        
        # For unknown extensions, try to read a small portion to check if it's text
        try:
            if file_path.stat().st_size > 10 * 1024 * 1024:  # Skip files larger than 10MB
                return False
            
            with open(file_path, 'rb') as f:
                chunk = f.read(1024)
                try:
                    chunk.decode('utf-8')
                    return True
                except UnicodeDecodeError:
                    return False
        except Exception:  # nosec B110 - Catch all file access errors
            return False
    
    async def _search_text_file(self, file_path: Path, pattern: re.Pattern, session_id: int) -> int:
        """Search for pattern in a text file"""
        try:
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                content = f.read()
            
            matches = list(pattern.finditer(content))
            if not matches:
                return 0
            
            # Create search result
            result = SearchResult(
                session_id=session_id,
                file_path=str(file_path).replace('/app/host_root', ''),  # Show original local path
                file_name=file_path.name,
                file_size=file_path.stat().st_size,
                file_type=file_path.suffix.lower(),
                match_count=len(matches),
                is_zip_file=False,
                preview_text=self._create_preview(content, matches[0].start())
            )
            self.db_session.add(result)
            self.db_session.flush()
            
            # Add match details
            lines = content.split('\n')
            for match in matches[:10]:  # Limit to first 10 matches per file
                line_num = content[:match.start()].count('\n') + 1
                line_content = lines[line_num - 1] if line_num <= len(lines) else ""
                
                match_detail = MatchDetail(
                    result_id=result.id,
                    line_number=line_num,
                    line_content=line_content,
                    match_position=match.start() - content.rfind('\n', 0, match.start()) - 1,
                    context_before='\n'.join(lines[max(0, line_num-3):line_num-1]),
                    context_after='\n'.join(lines[line_num:min(len(lines), line_num+2)])
                )
                self.db_session.add(match_detail)
            
            self.db_session.commit()
            return len(matches)
            
        except Exception as e:
            print(f"Error reading file {file_path}: {e}")
            return 0
    
    async def _search_zip_file(self, zip_path: Path, pattern: re.Pattern, session_id: int) -> int:
        """Search for pattern inside zip files"""
        total_matches = 0
        
        try:
            with zipfile.ZipFile(zip_path, 'r') as zip_file:
                for file_info in zip_file.filelist:
                    if file_info.is_dir():
                        continue
                    
                    file_name = Path(file_info.filename).name
                    file_ext = Path(file_info.filename).suffix.lower()
                    
                    if file_ext not in self.supported_text_extensions:
                        continue
                    
                    try:
                        with zip_file.open(file_info) as f:
                            content = f.read().decode('utf-8', errors='ignore')
                        
                        matches = list(pattern.finditer(content))
                        if matches:
                            result = SearchResult(
                                session_id=session_id,
                                file_path=file_info.filename,
                                file_name=file_name,
                                file_size=file_info.file_size,
                                file_type=file_ext,
                                match_count=len(matches),
                                is_zip_file=True,
                                zip_parent_path=str(zip_path).replace('/app/host_root', ''),  # Show original local path
                                preview_text=self._create_preview(content, matches[0].start())
                            )
                            self.db_session.add(result)
                            self.db_session.flush()
                            
                            # Add match details (limit to first 5 matches in zip files)
                            lines = content.split('\n')
                            for match in matches[:5]:
                                line_num = content[:match.start()].count('\n') + 1
                                line_content = lines[line_num - 1] if line_num <= len(lines) else ""
                                
                                match_detail = MatchDetail(
                                    result_id=result.id,
                                    line_number=line_num,
                                    line_content=line_content,
                                    match_position=match.start() - content.rfind('\n', 0, match.start()) - 1
                                )
                                self.db_session.add(match_detail)
                            
                            total_matches += len(matches)
                    
                    except Exception as e:
                        print(f"Error reading {file_info.filename} from {zip_path}: {e}")
                        continue
                
                self.db_session.commit()
                
        except Exception as e:
            print(f"Error reading zip file {zip_path}: {e}")
        
        return total_matches
    
    async def _search_tar_file(self, tar_path: Path, pattern: re.Pattern, session_id: int) -> int:
        """Search for pattern inside tar/tar.gz/tar.bz2 files"""
        total_matches = 0
        
        try:
            # Determine compression mode
            if tar_path.suffix in {'.gz', '.tgz'} or str(tar_path).endswith('.tar.gz'):
                mode = 'r:gz'
            elif tar_path.suffix in {'.bz2'} or str(tar_path).endswith('.tar.bz2'):
                mode = 'r:bz2'
            else:
                mode = 'r'
            
            with tarfile.open(tar_path, mode) as tar_file:
                for member in tar_file.getmembers():
                    if not member.isfile():
                        continue
                    
                    file_name = Path(member.name).name
                    file_ext = Path(member.name).suffix.lower()
                    
                    if file_ext not in self.supported_text_extensions:
                        continue
                    
                    try:
                        f = tar_file.extractfile(member)
                        if f is None:
                            continue
                        
                        content = f.read().decode('utf-8', errors='ignore')
                        f.close()
                        
                        matches = list(pattern.finditer(content))
                        if matches:
                            result = SearchResult(
                                session_id=session_id,
                                file_path=member.name,
                                file_name=file_name,
                                file_size=member.size,
                                file_type=file_ext,
                                match_count=len(matches),
                                is_zip_file=True,
                                zip_parent_path=str(tar_path).replace('/app/host_root', ''),
                                preview_text=self._create_preview(content, matches[0].start())
                            )
                            self.db_session.add(result)
                            self.db_session.flush()
                            
                            # Add match details (limit to first 5 matches)
                            lines = content.split('\n')
                            for match in matches[:5]:
                                line_num = content[:match.start()].count('\n') + 1
                                line_content = lines[line_num - 1] if line_num <= len(lines) else ""
                                
                                match_detail = MatchDetail(
                                    result_id=result.id,
                                    line_number=line_num,
                                    line_content=line_content,
                                    match_position=match.start() - content.rfind('\n', 0, match.start()) - 1
                                )
                                self.db_session.add(match_detail)
                            
                            total_matches += len(matches)
                    
                    except Exception as e:
                        print(f"Error reading {member.name} from {tar_path}: {e}")
                        continue
                
                self.db_session.commit()
                
        except Exception as e:
            print(f"Error reading tar file {tar_path}: {e}")
        
        return total_matches
    
    async def _search_docx_file(self, docx_path: Path, pattern: re.Pattern, session_id: int) -> int:
        """Search for pattern inside DOCX files"""
        if DocxDocument is None:
            return 0
        
        total_matches = 0
        try:
            doc = DocxDocument(str(docx_path))
            content = '\n'.join([paragraph.text for paragraph in doc.paragraphs])
            
            matches = list(pattern.finditer(content))
            if matches:
                result = SearchResult(
                    session_id=session_id,
                    file_path=str(docx_path).replace('/app/host_root', ''),
                    file_name=docx_path.name,
                    file_size=docx_path.stat().st_size,
                    file_type=docx_path.suffix.lower(),
                    match_count=len(matches),
                    is_zip_file=False,
                    preview_text=self._create_preview(content, matches[0].start())
                )
                self.db_session.add(result)
                self.db_session.flush()
                
                # Add match details
                for match in matches[:10]:
                    line_num = content[:match.start()].count('\n') + 1
                    lines = content.split('\n')
                    line_content = lines[line_num - 1] if line_num <= len(lines) else ""
                    
                    match_detail = MatchDetail(
                        result_id=result.id,
                        line_number=line_num,
                        line_content=line_content,
                        match_position=match.start() - content.rfind('\n', 0, match.start()) - 1
                    )
                    self.db_session.add(match_detail)
                
                total_matches = len(matches)
                self.db_session.commit()
                
        except Exception as e:
            print(f"Error reading DOCX file {docx_path}: {e}")
        
        return total_matches
    
    async def _search_xlsx_file(self, xlsx_path: Path, pattern: re.Pattern, session_id: int) -> int:
        """Search for pattern inside XLSX files"""
        if load_workbook is None:
            return 0
        
        total_matches = 0
        try:
            workbook = load_workbook(str(xlsx_path), read_only=True, data_only=True)
            content_parts = []
            
            for sheet in workbook.worksheets:
                for row in sheet.iter_rows(values_only=True):
                    row_text = ' '.join([str(cell) if cell is not None else '' for cell in row])
                    if row_text.strip():
                        content_parts.append(row_text)
            
            content = '\n'.join(content_parts)
            workbook.close()
            
            matches = list(pattern.finditer(content))
            if matches:
                result = SearchResult(
                    session_id=session_id,
                    file_path=str(xlsx_path).replace('/app/host_root', ''),
                    file_name=xlsx_path.name,
                    file_size=xlsx_path.stat().st_size,
                    file_type=xlsx_path.suffix.lower(),
                    match_count=len(matches),
                    is_zip_file=False,
                    preview_text=self._create_preview(content, matches[0].start())
                )
                self.db_session.add(result)
                self.db_session.flush()
                
                # Add match details
                for match in matches[:10]:
                    line_num = content[:match.start()].count('\n') + 1
                    lines = content.split('\n')
                    line_content = lines[line_num - 1] if line_num <= len(lines) else ""
                    
                    match_detail = MatchDetail(
                        result_id=result.id,
                        line_number=line_num,
                        line_content=line_content,
                        match_position=match.start() - content.rfind('\n', 0, match.start()) - 1
                    )
                    self.db_session.add(match_detail)
                
                total_matches = len(matches)
                self.db_session.commit()
                
        except Exception as e:
            print(f"Error reading XLSX file {xlsx_path}: {e}")
        
        return total_matches
    
    async def _search_pptx_file(self, pptx_path: Path, pattern: re.Pattern, session_id: int) -> int:
        """Search for pattern inside PPTX files"""
        if Presentation is None:
            return 0
        
        total_matches = 0
        try:
            prs = Presentation(str(pptx_path))
            content_parts = []
            
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        content_parts.append(shape.text)
            
            content = '\n'.join(content_parts)
            
            matches = list(pattern.finditer(content))
            if matches:
                result = SearchResult(
                    session_id=session_id,
                    file_path=str(pptx_path).replace('/app/host_root', ''),
                    file_name=pptx_path.name,
                    file_size=pptx_path.stat().st_size,
                    file_type=pptx_path.suffix.lower(),
                    match_count=len(matches),
                    is_zip_file=False,
                    preview_text=self._create_preview(content, matches[0].start())
                )
                self.db_session.add(result)
                self.db_session.flush()
                
                # Add match details
                for match in matches[:10]:
                    line_num = content[:match.start()].count('\n') + 1
                    lines = content.split('\n')
                    line_content = lines[line_num - 1] if line_num <= len(lines) else ""
                    
                    match_detail = MatchDetail(
                        result_id=result.id,
                        line_number=line_num,
                        line_content=line_content,
                        match_position=match.start() - content.rfind('\n', 0, match.start()) - 1
                    )
                    self.db_session.add(match_detail)
                
                total_matches = len(matches)
                self.db_session.commit()
                
        except Exception as e:
            print(f"Error reading PPTX file {pptx_path}: {e}")
        
        return total_matches
    
    async def _search_odf_file(self, odf_path: Path, pattern: re.Pattern, session_id: int) -> int:
        """Search for pattern inside ODF files (ODT, ODS, ODP)"""
        if odf_load is None or OdfP is None:
            return 0
        
        total_matches = 0
        try:
            doc = odf_load(str(odf_path))
            content_parts = []
            
            # Extract text from all paragraph elements
            for paragraph in doc.getElementsByType(OdfP):
                text = str(paragraph)
                if text.strip():
                    content_parts.append(text)
            
            content = '\n'.join(content_parts)
            
            matches = list(pattern.finditer(content))
            if matches:
                result = SearchResult(
                    session_id=session_id,
                    file_path=str(odf_path).replace('/app/host_root', ''),
                    file_name=odf_path.name,
                    file_size=odf_path.stat().st_size,
                    file_type=odf_path.suffix.lower(),
                    match_count=len(matches),
                    is_zip_file=False,
                    preview_text=self._create_preview(content, matches[0].start())
                )
                self.db_session.add(result)
                self.db_session.flush()
                
                # Add match details
                for match in matches[:10]:
                    line_num = content[:match.start()].count('\n') + 1
                    lines = content.split('\n')
                    line_content = lines[line_num - 1] if line_num <= len(lines) else ""
                    
                    match_detail = MatchDetail(
                        result_id=result.id,
                        line_number=line_num,
                        line_content=line_content,
                        match_position=match.start() - content.rfind('\n', 0, match.start()) - 1
                    )
                    self.db_session.add(match_detail)
                
                total_matches = len(matches)
                self.db_session.commit()
                
        except Exception as e:
            print(f"Error reading ODF file {odf_path}: {e}")
        
        return total_matches
    
    def _create_preview(self, content: str, match_position: int, context_length: int = 200) -> str:
        """Create a preview of the content around the match"""
        start = max(0, match_position - context_length // 2)
        end = min(len(content), match_position + context_length // 2)
        preview = content[start:end]
        
        if start > 0:
            preview = "..." + preview
        if end < len(content):
            preview = preview + "..."
        
        return preview

# WebSocket manager
class ConnectionManager:
    def __init__(self):
        self.active_connections: List[WebSocket] = []

    async def connect(self, websocket: WebSocket):
        await websocket.accept()
        self.active_connections.append(websocket)

    def disconnect(self, websocket: WebSocket):
        if websocket in self.active_connections:
            self.active_connections.remove(websocket)

    async def send_personal_message(self, message: str, websocket: WebSocket):
        await websocket.send_text(message)

manager = ConnectionManager()

# FastAPI app
app = FastAPI(title="File Search API", version="1.0.0")

app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],  # Configure appropriately for production
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

def get_db():
    db = SessionLocal()
    try:
        yield db
    finally:
        db.close()

@app.post("/search", response_model=dict)
async def start_search(request: SearchRequest):
    """Start a new search"""
    db = SessionLocal()
    try:
        search_engine = FileSearchEngine(db)
        session_id = await search_engine.search_files(request)
        return {"session_id": session_id, "status": "started"}
    finally:
        db.close()

@app.websocket("/ws/search")
async def websocket_search(websocket: WebSocket):
    """WebSocket endpoint for real-time search"""
    await manager.connect(websocket)
    try:
        # Wait for search request
        data = await websocket.receive_json()
        request = SearchRequest(**data)
        
        db = SessionLocal()
        try:
            search_engine = FileSearchEngine(db)
            session_id = await search_engine.search_files(request, websocket)
        finally:
            db.close()
            
    except WebSocketDisconnect:
        manager.disconnect(websocket)
    except Exception as e:
        print(f"ERROR in WebSocket search: {str(e)}")
        await websocket.send_json({"type": "error", "message": "An error occurred during the search operation. Please try again."})

@app.get("/sessions", response_model=List[SearchSessionResponse])
async def get_search_sessions():
    """Get all search sessions"""
    db = SessionLocal()
    try:
        sessions = db.query(SearchSession).order_by(SearchSession.start_time.desc()).all()
        return sessions
    finally:
        db.close()

@app.get("/sessions/{session_id}", response_model=SearchSessionResponse)
async def get_search_session(session_id: int):
    """Get a specific search session"""
    db = SessionLocal()
    try:
        session = db.query(SearchSession).filter(SearchSession.id == session_id).first()
        if not session:
            raise HTTPException(status_code=404, detail="Session not found")
        return session
    finally:
        db.close()

@app.get("/sessions/{session_id}/results", response_model=List[SearchResultResponse])
async def get_search_results(session_id: int, limit: int = 100, offset: int = 0):
    """Get search results for a session"""
    db = SessionLocal()
    try:
        results = db.query(SearchResult).filter(
            SearchResult.session_id == session_id
        ).offset(offset).limit(limit).all()
        return results
    finally:
        db.close()

@app.get("/results/{result_id}/matches", response_model=List[MatchDetailResponse])
async def get_match_details(result_id: int):
    """Get match details for a specific result"""
    db = SessionLocal()
    try:
        matches = db.query(MatchDetail).filter(MatchDetail.result_id == result_id).all()
        return matches
    finally:
        db.close()

@app.delete("/sessions/{session_id}")
async def delete_search_session(session_id: int):
    """Delete a search session and all its results"""
    db = SessionLocal()
    try:
        session = db.query(SearchSession).filter(SearchSession.id == session_id).first()
        if not session:
            raise HTTPException(status_code=404, detail="Session not found")
        
        db.delete(session)
        db.commit()
        return {"message": "Session deleted successfully"}
    finally:
        db.close()

@app.get("/results/{result_id}/preview")
async def get_file_preview(result_id: int, max_lines: int = 100):
    """Get file content preview for a specific result"""
    db = SessionLocal()
    try:
        result = db.query(SearchResult).filter(SearchResult.id == result_id).first()
        if not result:
            raise HTTPException(status_code=404, detail="Result not found")
        
        # Handle files inside archives
        if result.is_zip_file:
            return {
                "content": f"Preview not available for files inside archives.\n\nArchive: {result.zip_parent_path}\nFile: {result.file_path}\n\nTo view this file, extract the archive first.",
                "file_type": result.file_type,
                "is_binary": False,
                "truncated": False
            }
        
        # Construct full file path
        # The file_path in DB is already the local path (e.g., /Users/vndr/...)
        # We need to convert it to the Docker mounted path
        base_path = "/app/host_root"
        file_path_str = result.file_path
        
        # If path doesn't start with /, add it
        if not file_path_str.startswith('/'):
            file_path_str = '/' + file_path_str
        
        full_path = Path(base_path + file_path_str)
        
        print(f"DEBUG: Attempting to read file:")
        print(f"  - Result ID: {result_id}")
        print(f"  - DB file_path: {result.file_path}")
        print(f"  - Full path: {full_path}")
        print(f"  - Exists: {full_path.exists()}")
        
        if not full_path.exists():
            # Try alternative: maybe the path is already absolute in container
            alt_path = Path(file_path_str)
            print(f"  - Alternative path: {alt_path}")
            print(f"  - Alternative exists: {alt_path.exists()}")
            
            if alt_path.exists():
                full_path = alt_path
            else:
                return {
                    "content": f"File not found on disk.\n\nSearched paths:\n1. {full_path}\n2. {alt_path}\n\nThe file may have been moved or deleted since the search.\n\nOriginal path from search: {result.file_path}",
                    "file_type": result.file_type,
                    "is_binary": False,
                    "truncated": False
                }
        
        # Check if file is too large (limit to 10MB for preview)
        file_size = full_path.stat().st_size
        if file_size > 10 * 1024 * 1024:
            return {
                "content": f"File is too large to preview ({file_size:,} bytes).\n\nMaximum preview size: 10 MB",
                "file_type": result.file_type,
                "is_binary": False,
                "truncated": True,
                "file_size": file_size
            }
        
        file_ext = full_path.suffix.lower()
        
        # Handle different file types
        try:
            # Text files
            if file_ext in {'.txt', '.log', '.md', '.py', '.js', '.ts', '.html', '.css', 
                           '.json', '.xml', '.yml', '.yaml', '.ini', '.cfg', '.conf',
                           '.sh', '.bash', '.sql', '.csv', '.properties'}:
                with open(full_path, 'r', encoding='utf-8', errors='ignore') as f:
                    lines = f.readlines()
                    truncated = len(lines) > max_lines
                    content = ''.join(lines[:max_lines])
                    if truncated:
                        content += f"\n\n... (showing first {max_lines} of {len(lines)} lines)"
                    
                    return {
                        "content": content,
                        "file_type": file_ext,
                        "is_binary": False,
                        "truncated": truncated,
                        "total_lines": len(lines)
                    }
            
            # PDF files - extract text using PyPDF2
            elif file_ext == '.pdf' and PdfReader:
                try:
                    reader = PdfReader(str(full_path))
                    pages_text = []
                    total_pages = len(reader.pages)
                    
                    for page_num, page in enumerate(reader.pages, start=1):
                        try:
                            text = page.extract_text()
                            if text.strip():
                                pages_text.append(f"--- Page {page_num} ---\n{text}")
                        except Exception as e:
                            print(f"ERROR extracting PDF page {page_num} from {full_path.name}: {str(e)}")
                            pages_text.append(f"--- Page {page_num} ---\n[Unable to extract text from this page]")
                    
                    # Join all pages and split into lines for truncation
                    full_text = '\n\n'.join(pages_text)
                    lines = full_text.split('\n')
                    truncated = len(lines) > max_lines
                    content = '\n'.join(lines[:max_lines])
                    
                    if truncated:
                        content += f"\n\n... (showing first {max_lines} of {len(lines)} lines from {total_pages} pages)"
                    
                    return {
                        "content": content,
                        "file_type": file_ext,
                        "is_binary": False,
                        "truncated": truncated,
                        "total_pages": total_pages,
                        "total_lines": len(lines)
                    }
                except Exception as e:
                    print(f"ERROR reading PDF file {full_path.name}: {str(e)}")
                    return {
                        "content": f"Error reading PDF file.\n\nFilename: {full_path.name}\nSize: {file_size:,} bytes\n\n⚠️ Unable to extract text from this PDF file.",
                        "file_type": file_ext,
                        "is_binary": True,
                        "truncated": False
                    }
            
            # PDF files without PyPDF2 library
            elif file_ext == '.pdf':
                return {
                    "content": f"PDF Preview:\n\nFilename: {full_path.name}\nSize: {file_size:,} bytes\n\n⚠️ PyPDF2 library not installed.\nInstall it to enable PDF text extraction.",
                    "file_type": file_ext,
                    "is_binary": True,
                    "truncated": False
                }
            
            # DOCX files
            elif file_ext in {'.docx', '.doc'} and DocxDocument:
                try:
                    doc = DocxDocument(str(full_path))
                    paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
                    truncated = len(paragraphs) > max_lines
                    content = '\n'.join(paragraphs[:max_lines])
                    if truncated:
                        content += f"\n\n... (showing first {max_lines} of {len(paragraphs)} paragraphs)"
                    
                    return {
                        "content": content,
                        "file_type": file_ext,
                        "is_binary": False,
                        "truncated": truncated,
                        "total_paragraphs": len(paragraphs)
                    }
                except Exception as e:
                    print(f"ERROR reading DOCX file {full_path.name}: {str(e)}")
                    return {
                        "content": "Error reading DOCX file. The file may be corrupted or in an unsupported format.",
                        "file_type": file_ext,
                        "is_binary": True,
                        "truncated": False
                    }
            
            # XLSX files
            elif file_ext in {'.xlsx', '.xls'} and load_workbook:
                try:
                    workbook = load_workbook(str(full_path), read_only=True, data_only=True)
                    content_parts = []
                    
                    for sheet in workbook.worksheets[:5]:  # Limit to first 5 sheets
                        content_parts.append(f"\n=== Sheet: {sheet.title} ===\n")
                        row_count = 0
                        for row in sheet.iter_rows(values_only=True):
                            if row_count >= max_lines:
                                content_parts.append(f"\n... (showing first {max_lines} rows)")
                                break
                            row_text = '\t'.join([str(cell) if cell is not None else '' for cell in row])
                            if row_text.strip():
                                content_parts.append(row_text)
                                row_count += 1
                    
                    workbook.close()
                    
                    return {
                        "content": '\n'.join(content_parts),
                        "file_type": file_ext,
                        "is_binary": False,
                        "truncated": row_count >= max_lines
                    }
                except Exception as e:
                    print(f"ERROR reading Excel file {full_path.name}: {str(e)}")
                    return {
                        "content": "Error reading Excel file. The file may be corrupted or in an unsupported format.",
                        "file_type": file_ext,
                        "is_binary": True,
                        "truncated": False
                    }
            
            # PPTX files
            elif file_ext in {'.pptx', '.ppt'} and Presentation:
                try:
                    prs = Presentation(str(full_path))
                    content_parts = []
                    
                    for idx, slide in enumerate(prs.slides[:20], 1):  # Limit to first 20 slides
                        content_parts.append(f"\n=== Slide {idx} ===\n")
                        for shape in slide.shapes:
                            if hasattr(shape, "text") and shape.text.strip():
                                content_parts.append(shape.text)
                    
                    return {
                        "content": '\n'.join(content_parts),
                        "file_type": file_ext,
                        "is_binary": False,
                        "truncated": len(prs.slides) > 20,
                        "total_slides": len(prs.slides)
                    }
                except Exception as e:
                    print(f"ERROR reading PowerPoint file {full_path.name}: {str(e)}")
                    return {
                        "content": "Error reading PowerPoint file. The file may be corrupted or in an unsupported format.",
                        "file_type": file_ext,
                        "is_binary": True,
                        "truncated": False
                    }
            
            # Image files
            elif file_ext in {'.jpg', '.jpeg', '.png', '.gif', '.bmp', '.svg', '.ico'}:
                return {
                    "content": f"🖼️ Image File\n\nFilename: {full_path.name}\nType: {file_ext}\nSize: {file_size:,} bytes\n\n⚠️ Image preview not available in text mode.\nDownload the file to view.",
                    "file_type": file_ext,
                    "is_binary": True,
                    "truncated": False
                }
            
            # Binary files
            else:
                return {
                    "content": f"Binary file: {full_path.name}\n\nType: {file_ext}\nSize: {file_size:,} bytes\n\n⚠️ Binary file preview not available.\nThis appears to be a binary file type.",
                    "file_type": file_ext,
                    "is_binary": True,
                    "truncated": False
                }
                
        except Exception as e:
            print(f"ERROR reading file: {str(e)}")
            import traceback
            traceback.print_exc()
            return {
                "content": "Error reading file. The file may be inaccessible, corrupted, or in an unsupported format.",
                "file_type": result.file_type,
                "is_binary": False,
                "truncated": False
            }
    
    except Exception as e:
        print(f"ERROR in preview endpoint: {str(e)}")
        import traceback
        traceback.print_exc()
        db.close()
        # Return error as content instead of raising exception
        return {
            "content": "Error loading preview. Unable to load the file preview. Please check if the file exists and is accessible.",
            "file_type": "error",
            "is_binary": False,
            "truncated": False
        }
    
    finally:
        db.close()

@app.get("/filesystem/directories")
async def list_directories(path: str = "/"):
    """
    List directories at the specified path for directory browser.
    Returns subdirectories only (not files).
    """
    try:
        # Security: Validate and resolve path securely
        host_root = Path("/app/host_root")
        validated_path = validate_and_resolve_path(path, host_root)
        
        # Check if path exists and is a directory
        if not os.path.exists(validated_path):
            raise HTTPException(status_code=404, detail="Path not found")
        
        if not os.path.isdir(validated_path):
            raise HTTPException(status_code=400, detail="Path is not a directory")
        
        # List subdirectories
        directories = []
        try:
            for item_name in sorted(os.listdir(validated_path)):
                item_path = os.path.join(validated_path, item_name)
                if os.path.isdir(item_path):
                    # Check for read permission
                    if os.access(item_path, os.R_OK):
                        directories.append({
                            "name": item_name,
                            "path": item_path,
                            "display_path": item_path.replace("/app/host_root", "")
                        })
        except PermissionError:
            pass  # Skip directories we can't read
        
        # Get parent directory info
        host_root_resolved = str(host_root.resolve())
        parent_path = None
        if validated_path != host_root_resolved:
            parent = os.path.dirname(validated_path)
            if parent.startswith(host_root_resolved):
                parent_path = parent.replace("/app/host_root", "")
        
        return {
            "current_path": validated_path.replace("/app/host_root", ""),
            "parent_path": parent_path,
            "directories": directories
        }
        
    except HTTPException:
        raise
    except Exception as e:
        print(f"ERROR listing directories: {str(e)}")
        import traceback
        traceback.print_exc()
        raise HTTPException(status_code=500, detail="Error listing directories. Unable to access the specified path.")

@app.post("/api/analyze-directory")
async def analyze_directory(path: str, find_duplicates: bool = True, max_hash_size: int = 10, session_id: str = None, min_size: int = None, max_size: int = None):
    """
    Analyze a directory and return comprehensive statistics about its contents.
    Returns file type distribution, size statistics, duplicate files, and detailed file list.
    
    Parameters:
    - path: Directory path to analyze
    - find_duplicates: Whether to calculate file hashes for duplicate detection (slower)
    - max_hash_size: Maximum file size in MB to hash for duplicates (default: 10MB)
    - session_id: Optional session ID for cancellation support
    - min_size: Optional minimum file size in bytes to include in analysis
    - max_size: Optional maximum file size in bytes to include in analysis
    """
    import hashlib
    from collections import defaultdict
    from concurrent.futures import ThreadPoolExecutor, as_completed
    import threading
    import uuid
    
    # Generate session ID if not provided
    if not session_id:
        session_id = str(uuid.uuid4())
    
    # Register this analysis session
    analysis_cancellations[session_id] = False
    
    def is_cancelled():
        """Check if analysis has been cancelled"""
        return analysis_cancellations.get(session_id, False)
    
    def quick_hash(file_path: Path, file_size: int) -> str:
        """
        Fast hashing using size-based sampling:
        - Files < 1MB: Full hash
        - Files >= 1MB: Hash first 64KB + middle 64KB + last 64KB + size
        """
        try:
            # MD5 is used for file content comparison (duplicate detection), not security
            hash_md5 = hashlib.md5(usedforsecurity=False)
            
            # Always include file size in hash
            hash_md5.update(str(file_size).encode())
            
            if file_size < 1024 * 1024:  # < 1MB: full hash
                with open(file_path, "rb") as f:
                    for chunk in iter(lambda: f.read(8192), b""):
                        hash_md5.update(chunk)
            else:  # >= 1MB: sample hash
                with open(file_path, "rb") as f:
                    # First 64KB
                    hash_md5.update(f.read(65536))
                    
                    # Middle 64KB
                    if file_size > 131072:
                        f.seek(file_size // 2)
                        hash_md5.update(f.read(65536))
                    
                    # Last 64KB
                    if file_size > 65536:
                        f.seek(-65536, 2)
                        hash_md5.update(f.read(65536))
            
            return hash_md5.hexdigest()
        except:
            return None
    
    try:
        # Security: Validate and resolve path securely
        host_root = Path("/app/host_root")
        validated_path = validate_and_resolve_path(path, host_root)
        
        # Check if path exists and is a directory
        if not os.path.exists(validated_path):
            raise HTTPException(status_code=404, detail="Path not found")
        
        if not os.path.isdir(validated_path):
            raise HTTPException(status_code=400, detail="Path is not a directory")
        
        # Initialize data structures
        file_list = []
        file_type_counts = defaultdict(int)
        file_type_sizes = defaultdict(int)
        size_key_map = {}  # Map size to list of files for pre-filtering duplicates
        total_size = 0
        total_files = 0
        total_dirs = 0
        empty_dirs = []  # Track empty directories
        
        max_hash_bytes = max_hash_size * 1024 * 1024
        
        print(f"Starting directory analysis: {validated_path}")
        print(f"Find duplicates: {find_duplicates}, Max hash size: {max_hash_size}MB")
        print(f"Excluding hidden directories (starting with '.')")
        
        # Phase 1: Fast metadata collection (no hashing)
        print("Phase 1: Collecting file metadata...")
        for root, dirs, files in os.walk(validated_path):
            # Check for cancellation
            if is_cancelled():
                print(f"Analysis cancelled by user at {total_files} files")
                break
            
            # Filter out hidden directories (starting with ".")
            dirs[:] = [d for d in dirs if not d.startswith('.')]
            total_dirs += len(dirs)
            
            # Check for empty directories (no files and no subdirectories)
            if len(files) == 0 and len(dirs) == 0 and root != validated_path:
                # This is an empty directory
                empty_dir_path = str(Path(root)).replace("/app/host_root", "")
                empty_dirs.append({
                    "path": empty_dir_path,
                    "name": Path(root).name,
                    "parent": str(Path(root).parent).replace("/app/host_root", "")
                })
            
            for filename in files:
                # Check cancellation periodically (every 100 files)
                if total_files % 100 == 0 and is_cancelled():
                    print(f"Analysis cancelled by user at {total_files} files")
                    break
                
                try:
                    file_path = os.path.join(root, filename)
                    
                    # Check if we can access the file
                    if not os.access(file_path, os.R_OK):
                        continue
                    
                    # Get file stats (fast - just metadata)
                    stat = os.stat(file_path)
                    file_size = stat.st_size
                    modified_time = datetime.fromtimestamp(stat.st_mtime)
                    
                    # Apply size filters if specified
                    if min_size is not None and file_size < min_size:
                        continue
                    if max_size is not None and file_size > max_size:
                        continue
                    
                    # Get file extension and MIME type
                    file_ext = file_path.suffix.lower() or '.no_extension'
                    mime_type, _ = mimetypes.guess_type(str(file_path))
                    mime_type = mime_type or 'application/octet-stream'
                    
                    # Add to file list (without hash yet)
                    file_info = {
                        "name": filename,
                        "path": str(file_path).replace("/app/host_root", ""),
                        "size": file_size,
                        "extension": file_ext,
                        "mime_type": mime_type,
                        "modified": modified_time.isoformat(),
                        "hash": None,
                        "_full_path": str(file_path)  # Temporary for hashing
                    }
                    file_list.append(file_info)
                    
                    # Update statistics
                    file_type_counts[file_ext] += 1
                    file_type_sizes[file_ext] += file_size
                    total_size += file_size
                    total_files += 1
                    
                    # Pre-filter for duplicates: only hash files with same size
                    if find_duplicates and file_size > 0 and file_size <= max_hash_bytes:
                        if file_size not in size_key_map:
                            size_key_map[file_size] = []
                        size_key_map[file_size].append(file_info)
                
                except (PermissionError, OSError):
                    continue  # Skip files we can't access
        
        print(f"Phase 1 complete: {total_files} files, {total_dirs} directories")
        
        # Phase 2: Parallel hashing (only for potential duplicates)
        file_hashes = defaultdict(list)
        
        if find_duplicates:
            # Only hash files that have at least one other file with the same size
            files_to_hash = []
            for size, file_infos in size_key_map.items():
                if len(file_infos) > 1:  # Only hash if multiple files have same size
                    files_to_hash.extend(file_infos)
            
            if files_to_hash:
                print(f"Phase 2: Hashing {len(files_to_hash)} potential duplicates (parallel)...")
                
                # Use thread pool for parallel hashing
                with ThreadPoolExecutor(max_workers=min(8, os.cpu_count() or 4)) as executor:
                    future_to_file = {
                        executor.submit(quick_hash, Path(f["_full_path"]), f["size"]): f
                        for f in files_to_hash
                    }
                    
                    completed = 0
                    for future in as_completed(future_to_file):
                        file_info = future_to_file[future]
                        try:
                            file_hash = future.result()
                            if file_hash:
                                file_info["hash"] = file_hash
                                file_hashes[file_hash].append(file_info)
                            completed += 1
                            if completed % 1000 == 0:
                                print(f"  Hashed {completed}/{len(files_to_hash)} files...")
                        except Exception as e:
                            print(f"  Error hashing {file_info['name']}: {e}")
                            continue
                
                print(f"Phase 2 complete: Hashed {len(files_to_hash)} files")
        
        # Clean up temporary field
        for file_info in file_list:
            if "_full_path" in file_info:
                del file_info["_full_path"]
        
        print("Building final results...")
        
        # Track duplicates
        duplicates = []
        duplicate_size = 0
        if find_duplicates:
            for hash_value, files in file_hashes.items():
                if len(files) > 1:
                    # This is a duplicate group
                    group_size = files[0]["size"]
                    wasted_space = group_size * (len(files) - 1)
                    duplicate_size += wasted_space
                    
                    duplicates.append({
                        "hash": hash_value,
                        "count": len(files),
                        "size": group_size,
                        "wasted_space": wasted_space,
                        "files": [{"name": f["name"], "path": f["path"]} for f in files]
                    })
            
            # Sort duplicates by wasted space
            duplicates.sort(key=lambda x: x["wasted_space"], reverse=True)
        
        # Convert file type stats to list format
        file_types = [
            {
                "extension": ext,
                "count": count,
                "total_size": file_type_sizes[ext],
                "average_size": file_type_sizes[ext] // count if count > 0 else 0
            }
            for ext, count in file_type_counts.items()
        ]
        
        # Sort by total size
        file_types.sort(key=lambda x: x["total_size"], reverse=True)
        
        # Calculate size distribution by ranges
        size_ranges = {
            "0-1KB": 0,
            "1KB-10KB": 0,
            "10KB-100KB": 0,
            "100KB-1MB": 0,
            "1MB-10MB": 0,
            "10MB-100MB": 0,
            "100MB+": 0
        }
        
        for file_info in file_list:
            size = file_info["size"]
            if size < 1024:
                size_ranges["0-1KB"] += 1
            elif size < 10 * 1024:
                size_ranges["1KB-10KB"] += 1
            elif size < 100 * 1024:
                size_ranges["10KB-100KB"] += 1
            elif size < 1024 * 1024:
                size_ranges["100KB-1MB"] += 1
            elif size < 10 * 1024 * 1024:
                size_ranges["1MB-10MB"] += 1
            elif size < 100 * 1024 * 1024:
                size_ranges["10MB-100MB"] += 1
            else:
                size_ranges["100MB+"] += 1
        
        # Get top 10 largest files
        largest_files = sorted(file_list, key=lambda x: x["size"], reverse=True)[:10]
        
        # Check if analysis was cancelled
        cancelled = is_cancelled()
        
        return {
            "session_id": session_id,
            "cancelled": cancelled,
            "summary": {
                "path": validated_path.replace("/app/host_root", ""),
                "total_files": total_files,
                "total_directories": total_dirs,
                "total_size": total_size,
                "duplicate_files": sum(d["count"] - 1 for d in duplicates),
                "wasted_space": duplicate_size,
                "unique_file_types": len(file_type_counts),
                "empty_directories": len(empty_dirs)
            },
            "file_types": file_types,
            "size_distribution": size_ranges,
            "duplicates": duplicates,
            "empty_directories": empty_dirs,
            "largest_files": largest_files,
            "all_files": file_list
        }
    
    except HTTPException:
        raise
    except Exception as e:
        print(f"ERROR analyzing directory: {str(e)}")
        import traceback
        traceback.print_exc()
        raise HTTPException(status_code=500, detail="Error analyzing directory. Unable to complete the analysis operation.")
    finally:
        # Clean up cancellation tracking
        if session_id in analysis_cancellations:
            del analysis_cancellations[session_id]

@app.post("/api/export-analysis-csv")
async def export_analysis_csv(data: dict, export_type: str = "all_files"):
    """
    Export analysis data to CSV format
    
    Parameters:
    - data: The analysis data to export
    - export_type: Type of export - 'all_files', 'duplicates', or 'file_types'
    """
    try:
        output = io.StringIO()
        
        if export_type == "all_files":
            # Export all files list
            writer = csv.writer(output)
            writer.writerow(['File Name', 'Path', 'Size (Bytes)', 'Size (Human)', 'Extension', 'MIME Type', 'Modified Date'])
            
            for file_info in data.get('all_files', []):
                size_bytes = file_info.get('size', 0)
                # Human readable size
                if size_bytes == 0:
                    size_human = '0 Bytes'
                else:
                    k = 1024
                    sizes = ['Bytes', 'KB', 'MB', 'GB', 'TB']
                    i = int(min(len(sizes) - 1, int((len(str(size_bytes)) - 1) / 3)))
                    size_human = f"{round(size_bytes / (k ** i), 2)} {sizes[i]}"
                
                writer.writerow([
                    file_info.get('name', ''),
                    file_info.get('path', ''),
                    size_bytes,
                    size_human,
                    file_info.get('extension', ''),
                    file_info.get('mime_type', ''),
                    file_info.get('modified', '')
                ])
            
            filename = f"file_list_{datetime.now().strftime('%Y%m%d_%H%M%S')}.csv"
        
        elif export_type == "duplicates":
            # Export duplicate files
            writer = csv.writer(output)
            writer.writerow(['Group', 'Duplicate Count', 'File Size (Bytes)', 'Wasted Space (Bytes)', 'File Name', 'File Path'])
            
            for idx, dup_group in enumerate(data.get('duplicates', []), 1):
                for file_info in dup_group.get('files', []):
                    writer.writerow([
                        f"Group {idx}",
                        dup_group.get('count', 0),
                        dup_group.get('size', 0),
                        dup_group.get('wasted_space', 0),
                        file_info.get('name', ''),
                        file_info.get('path', '')
                    ])
            
            filename = f"duplicates_{datetime.now().strftime('%Y%m%d_%H%M%S')}.csv"
        
        elif export_type == "file_types":
            # Export file type statistics
            writer = csv.writer(output)
            writer.writerow(['Extension', 'File Count', 'Total Size (Bytes)', 'Total Size (Human)', 'Average Size (Bytes)'])
            
            for file_type in data.get('file_types', []):
                total_size = file_type.get('total_size', 0)
                # Human readable size
                if total_size == 0:
                    size_human = '0 Bytes'
                else:
                    k = 1024
                    sizes = ['Bytes', 'KB', 'MB', 'GB', 'TB']
                    i = int(min(len(sizes) - 1, int((len(str(total_size)) - 1) / 3)))
                    size_human = f"{round(total_size / (k ** i), 2)} {sizes[i]}"
                
                writer.writerow([
                    file_type.get('extension', ''),
                    file_type.get('count', 0),
                    total_size,
                    size_human,
                    file_type.get('average_size', 0)
                ])
            
            filename = f"file_types_{datetime.now().strftime('%Y%m%d_%H%M%S')}.csv"
        
        else:
            raise HTTPException(status_code=400, detail="Invalid export type")
        
        # Prepare response
        output.seek(0)
        return StreamingResponse(
            iter([output.getvalue()]),
            media_type="text/csv",
            headers={"Content-Disposition": f"attachment; filename={filename}"}
        )
    
    except Exception as e:
        print(f"ERROR exporting CSV: {str(e)}")
        import traceback
        traceback.print_exc()
        raise HTTPException(status_code=500, detail="Error exporting CSV. Unable to generate the export file.")

@app.post("/api/cancel-analysis/{session_id}")
async def cancel_analysis(session_id: str):
    """
    Cancel an ongoing directory analysis
    """
    if session_id in analysis_cancellations:
        analysis_cancellations[session_id] = True
        return {"status": "cancelled", "session_id": session_id}
    else:
        raise HTTPException(status_code=404, detail="Analysis session not found or already completed")

# Saved Searches Endpoints
@app.get("/api/saved-searches", response_model=List[SavedSearchResponse])
async def get_saved_searches():
    """
    Get all saved searches
    """
    db = SessionLocal()
    try:
        searches = db.query(SavedSearch).order_by(SavedSearch.created_at.desc()).all()
        return searches
    finally:
        db.close()

@app.post("/api/saved-searches", response_model=SavedSearchResponse)
async def create_saved_search(search: SavedSearchCreate):
    """
    Create a new saved search
    """
    db = SessionLocal()
    try:
        # Check if name already exists
        existing = db.query(SavedSearch).filter(SavedSearch.name == search.name).first()
        if existing:
            raise HTTPException(status_code=400, detail="A saved search with this name already exists")
        
        db_search = SavedSearch(
            name=search.name,
            search_term=search.search_term,
            search_path=search.search_path,
            case_sensitive=search.case_sensitive,
            include_zip_files=search.include_zip_files,
            search_filenames=search.search_filenames,
            file_type=search.file_type,
            min_size=search.min_size,
            max_size=search.max_size
        )
        db.add(db_search)
        db.commit()
        db.refresh(db_search)
        return db_search
    except HTTPException:
        raise
    except Exception as e:
        db.rollback()
        print(f"ERROR creating saved search: {str(e)}")
        raise HTTPException(status_code=500, detail="Error creating saved search. Unable to save the search configuration.")
    finally:
        db.close()

@app.delete("/api/saved-searches/{search_id}")
async def delete_saved_search(search_id: int):
    """
    Delete a saved search
    """
    db = SessionLocal()
    try:
        search = db.query(SavedSearch).filter(SavedSearch.id == search_id).first()
        if not search:
            raise HTTPException(status_code=404, detail="Saved search not found")
        
        db.delete(search)
        db.commit()
        return {"status": "success", "message": "Saved search deleted"}
    except HTTPException:
        raise
    except Exception as e:
        db.rollback()
        print(f"ERROR deleting saved search: {str(e)}")
        raise HTTPException(status_code=500, detail="Error deleting saved search. Unable to remove the search configuration.")
    finally:
        db.close()

@app.post("/api/delete-empty-directories")
async def delete_empty_directories(request: dict):
    """
    Delete empty directories safely
    
    Request body:
    {
        "directories": ["/path/to/dir1", "/path/to/dir2", ...]
    }
    
    Returns:
    {
        "deleted": ["/path/to/dir1", ...],
        "failed": [{"path": "/path/to/dir2", "error": "..."}, ...],
        "total_deleted": 3,
        "total_failed": 1
    }
    """
    import shutil
    
    directories = request.get("directories", [])
    print(f"Received request to delete {len(directories)} directories")
    print(f"Directories: {directories}")
    
    if not directories:
        raise HTTPException(status_code=400, detail="No directories provided")
    
    deleted = []
    failed = []
    host_root = Path("/app/host_root")
    
    for dir_path in directories:
        try:
            print(f"Processing directory: {dir_path}")
            # Security: Validate and resolve path securely
            validated_path = validate_and_resolve_path(dir_path, host_root)
            
            print(f"Validated path: {validated_path}")
            
            # Check if path exists and is a directory
            if not os.path.exists(validated_path):
                print(f"FAILED: Directory does not exist")
                failed.append({
                    "path": dir_path,
                    "error": "Directory does not exist"
                })
                continue
            
            if not os.path.isdir(validated_path):
                print(f"FAILED: Path is not a directory")
                failed.append({
                    "path": dir_path,
                    "error": "Path is not a directory"
                })
                continue
            
            # Double-check the directory is actually empty
            dir_contents = os.listdir(validated_path)
            if dir_contents:
                print(f"FAILED: Directory is not empty, contains: {dir_contents}")
                failed.append({
                    "path": dir_path,
                    "error": "Directory is not empty"
                })
                continue
            
            # Delete the empty directory
            os.rmdir(validated_path)
            deleted.append(dir_path)
            print(f"SUCCESS: Deleted empty directory: {dir_path}")
            
        except PermissionError:
            failed.append({
                "path": dir_path,
                "error": "Permission denied"
            })
        except Exception as e:
            failed.append({
                "path": dir_path,
                "error": "Unable to delete directory"
            })
            print(f"EXCEPTION deleting {dir_path}: {e}")
    
    result = {
        "deleted": deleted,
        "failed": failed,
        "total_deleted": len(deleted),
        "total_failed": len(failed)
    }
    print(f"Delete operation complete: {result}")
    return result

@app.get("/health")
async def health_check():
    """Health check endpoint"""
    return {"status": "healthy", "timestamp": datetime.utcnow()}

if __name__ == "__main__":
    # Create tables
    Base.metadata.create_all(bind=engine)
    
    # Run the application
    uvicorn.run(app, host="0.0.0.0", port=8000)  # nosec B104 - Binding to all interfaces is required for Docker